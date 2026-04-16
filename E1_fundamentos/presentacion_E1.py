#!/usr/bin/env python3
"""
Genera la presentación del Encuentro 1 de la capacitación ISELIN.
Estilo DATAOILERS: fondo oscuro, acentos azul/teal, tipografía limpia.

Uso: python3 presentacion_E1.py
Genera: presentacion_E1.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Constantes de diseño ───
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colores
DARK_BG = RGBColor(0x0F, 0x1B, 0x2D)
DARK_BG2 = RGBColor(0x13, 0x21, 0x37)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_TEAL = RGBColor(0x4A, 0xAD, 0xA4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x8F, 0xA8, 0xBF)
MID_GRAY = RGBColor(0x84, 0x94, 0xA7)
SUBTLE_GRAY = RGBColor(0x5A, 0x7A, 0x94)
CARD_BG = RGBColor(0x1A, 0x2A, 0x3A)
CARD_BG2 = RGBColor(0x16, 0x23, 0x3E)
TABLE_HEADER = RGBColor(0x1A, 0x2A, 0x3A)
TABLE_ROW_EVEN = RGBColor(0x11, 0x1D, 0x30)
TABLE_ROW_ODD = RGBColor(0x0F, 0x1B, 0x2D)
ORANGE_ACCENT = RGBColor(0xE6, 0x7E, 0x22)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)


# ─── Helpers ───

def set_slide_bg(slide, color):
    """Establece el color de fondo de un slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """Agrega una forma con relleno sólido."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Agrega un textbox con texto formateado."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(6),
                  font_name="Calibri"):
    """Agrega un párrafo a un text_frame existente."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Agrega una línea de acento horizontal."""
    return add_shape(slide, left, top, width, Pt(4), color)


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, bullet_color=ACCENT_TEAL):
    """Agrega una lista con bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        # Bullet level
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        # Remove existing bullets
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)
        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {
            'val': str(bullet_color)
        })
        buClr.append(srgbClr)
        for existing in pPr.findall(qn('a:buClr')):
            pPr.remove(existing)
        pPr.append(buClr)

    return txBox


def add_card(slide, left, top, width, height, title, body, bg=CARD_BG,
             accent=ACCENT_BLUE, title_size=16, body_size=13):
    """Agrega una tarjeta con título y cuerpo."""
    # Fondo de tarjeta con bordes redondeados
    card = add_shape(slide, left, top, width, height, bg,
                     MSO_SHAPE.ROUNDED_RECTANGLE)

    # Línea de acento arriba
    add_shape(slide, left + Inches(0.15), top + Inches(0.12),
              Inches(0.5), Pt(3), accent)

    # Título
    add_textbox(slide, left + Inches(0.2), top + Inches(0.25),
                width - Inches(0.4), Inches(0.4),
                title, font_size=title_size, color=ACCENT_TEAL, bold=True)

    # Cuerpo
    add_textbox(slide, left + Inches(0.2), top + Inches(0.65),
                width - Inches(0.4), height - Inches(0.8),
                body, font_size=body_size, color=LIGHT_GRAY)

    return card


def add_number_badge(slide, left, top, number, color=ACCENT_BLUE):
    """Agrega un número dentro de un círculo."""
    circle = add_shape(slide, left, top, Inches(0.55), Inches(0.55),
                       color, MSO_SHAPE.OVAL)
    add_textbox(slide, left, top + Inches(0.05), Inches(0.55), Inches(0.45),
                str(number), font_size=20, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    return circle


# ─── Slides ───

def slide_portada(prs):
    """Slide 1: Portada"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Barra lateral izquierda
    add_shape(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, ACCENT_BLUE)

    # Logo text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "data "
    run1.font.size = Pt(22)
    run1.font.color.rgb = MID_GRAY
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "oilers"
    run2.font.size = Pt(22)
    run2.font.color.rgb = ACCENT_TEAL
    run2.font.bold = True
    run2.font.name = "Calibri"

    # Línea de acento
    add_accent_line(slide, Inches(0.6), Inches(3.2), Inches(1.2))

    # Título principal
    add_textbox(slide, Inches(0.6), Inches(3.45), Inches(9), Inches(1.5),
                "Capacitación en Análisis de\nDatos con Power BI",
                font_size=40, color=WHITE, bold=True)

    # Subtítulo
    add_textbox(slide, Inches(0.6), Inches(4.95), Inches(8), Inches(0.8),
                "Encuentro 1: Fundamentos, Calidad y Primer Contacto con PBI",
                font_size=20, color=LIGHT_GRAY)

    # Info
    add_textbox(slide, Inches(0.6), Inches(6.3), Inches(6), Inches(0.5),
                "ISELIN  ·  Jueves 16 de abril de 2026  ·  Presencial  ·  2 horas",
                font_size=13, color=SUBTLE_GRAY)

    # Decoración derecha: rectángulo sutil
    rect = add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(4.5), CARD_BG2)
    rect.fill.fore_color.rgb = RGBColor(0x14, 0x20, 0x35)

    # Ícono PBI simulado (texto grande)
    add_textbox(slide, Inches(10.7), Inches(3.0), Inches(2.2), Inches(1.5),
                "PBI", font_size=60, color=RGBColor(0x1E, 0x30, 0x50),
                bold=True, alignment=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """Slide 2: Agenda del día"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Header
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                "Agenda del día", font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    # Bloques de agenda
    bloques = [
        ("10 min", "Apertura y contexto", ACCENT_BLUE),
        ("25 min", "Análisis de datos y calidad", ACCENT_TEAL),
        ("25 min", "Conceptos fundamentales", ACCENT_BLUE),
        ("35 min", "Power BI: demo en vivo", ACCENT_TEAL),
        ("15 min", "Cierre y tarea", ACCENT_BLUE),
        ("10 min", "Tiempo extra", MID_GRAY),
    ]

    start_y = Inches(1.5)
    for i, (tiempo, tema, color) in enumerate(bloques):
        y = start_y + Inches(i * 0.88)

        # Barra de color
        add_shape(slide, Inches(0.8), y, Pt(5), Inches(0.65), color)

        # Número
        add_textbox(slide, Inches(1.1), y + Inches(0.08), Inches(0.5), Inches(0.5),
                    str(i + 1) if i < 5 else "—",
                    font_size=22, color=color, bold=True)

        # Tema
        add_textbox(slide, Inches(1.7), y + Inches(0.05), Inches(5), Inches(0.35),
                    tema, font_size=20, color=WHITE, bold=True)

        # Duración
        add_textbox(slide, Inches(1.7), y + Inches(0.38), Inches(3), Inches(0.3),
                    tiempo, font_size=13, color=SUBTLE_GRAY)

    # Lado derecho: info del encuentro
    add_card(slide, Inches(8.5), Inches(1.5), Inches(4.2), Inches(2.2),
             "Encuentro 1 de 4", "Jueves 16/04/2026\nModalidad: Presencial\nDuración: 2 horas\n\nNivel: desde cero",
             title_size=14, body_size=14)

    # Objetivo
    add_card(slide, Inches(8.5), Inches(4.1), Inches(4.2), Inches(2.8),
             "Objetivo específico",
             "Establecer el \"para qué\" del análisis de datos, introducir los conceptos "
             "fundamentales de Power BI y lograr la primera carga de datos real.",
             accent=ACCENT_TEAL, title_size=14, body_size=13)


def slide_por_que_importa(prs):
    """Slide 3: ¿Por qué análisis de datos?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "¿Por qué importa el análisis de datos?",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                "Conectar el análisis de datos con la cultura de calidad de ISELIN",
                font_size=18, color=LIGHT_GRAY)

    # 3 tarjetas
    cards = [
        ("Datos → Información → Decisiones",
         "Los datos por sí solos no dicen nada.\nCuando los organizamos y visualizamos,\nse convierten en información para\ntomar mejores decisiones."),
        ("Cultura de calidad",
         "ISELIN ya trabaja con mejora continua.\nPower BI se inserta en ese proceso\ncomo la herramienta para verificar\nresultados con datos reales."),
        ("Visibilidad para todos",
         "Un dashboard permite que todas\nlas áreas vean los mismos números,\nalineando criterios y facilitando\nla comunicación."),
    ]

    for i, (title, body) in enumerate(cards):
        x = Inches(0.8) + Inches(i * 4.15)
        add_card(slide, x, Inches(2.2), Inches(3.8), Inches(3.5),
                 title, body, accent=[ACCENT_BLUE, ACCENT_TEAL, ACCENT_BLUE][i],
                 title_size=16, body_size=14)


def slide_ciclo_mejora(prs):
    """Slide 4: Ciclo de mejora continua (PDCA)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Ciclo de mejora continua",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.6),
                "Power BI es la herramienta de \"Verificar\" dentro del ciclo de mejora",
                font_size=18, color=LIGHT_GRAY)

    # 4 etapas del ciclo
    etapas = [
        ("1", "Planificar", "Definir objetivos\ny planificar acciones", ACCENT_BLUE),
        ("2", "Hacer", "Ejecutar lo\nplanificado", ACCENT_TEAL),
        ("3", "Verificar", "Medir resultados\ncon datos reales", ORANGE_ACCENT),
        ("4", "Actuar", "Ajustar y mejorar\nel proceso", GREEN_ACCENT),
    ]

    for i, (num, titulo, desc, color) in enumerate(etapas):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(2.5)

        # Número circular
        add_number_badge(slide, x + Inches(1.05), y, num, color)

        # Título etapa
        add_textbox(slide, x, y + Inches(0.75), Inches(2.8), Inches(0.5),
                    titulo, font_size=22, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Descripción
        add_textbox(slide, x, y + Inches(1.25), Inches(2.8), Inches(1.2),
                    desc, font_size=14, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Highlight en Verificar
    highlight = add_shape(slide, Inches(7.1), Inches(2.3),
                          Inches(2.85), Inches(3.0),
                          RGBColor(0xE6, 0x7E, 0x22))
    highlight.fill.fore_color.rgb = RGBColor(0x2A, 0x1A, 0x0A)
    # Borde sutil

    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(0.6),
                "▸  Power BI nos permite verificar con datos si las acciones que tomamos están dando resultado",
                font_size=16, color=ACCENT_TEAL)


def slide_etl(prs):
    """Slide 5: Proceso de carga (ETL)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Proceso de carga: del Excel a Power BI",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    # 3 pasos ETL
    pasos = [
        ("Extraer", "Sacar los datos\ndesde Excel, Wara\nu otra fuente", "📥"),
        ("Transformar", "Limpiar, corregir\nformatos, eliminar\ndatos innecesarios", "🔧"),
        ("Cargar", "Meter los datos\nlimpios en Power BI\nlistos para usar", "📊"),
    ]

    for i, (titulo, desc, icono) in enumerate(pasos):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(1.8)

        # Tarjeta
        card = add_shape(slide, x, y, Inches(3.8), Inches(3.2), CARD_BG,
                         MSO_SHAPE.ROUNDED_RECTANGLE)

        # Número de paso
        add_number_badge(slide, x + Inches(0.2), y + Inches(0.2),
                         str(i + 1),
                         [ACCENT_BLUE, ACCENT_TEAL, GREEN_ACCENT][i])

        # Título
        add_textbox(slide, x + Inches(0.9), y + Inches(0.25),
                    Inches(2.5), Inches(0.4),
                    titulo, font_size=22, color=WHITE, bold=True)

        # Descripción
        add_textbox(slide, x + Inches(0.3), y + Inches(0.9),
                    Inches(3.2), Inches(2.0),
                    desc, font_size=16, color=LIGHT_GRAY)

    # Flecha entre pasos (texto)
    for i in range(2):
        x = Inches(4.4) + Inches(i * 4.15)
        add_textbox(slide, x, Inches(3.1), Inches(0.5), Inches(0.5),
                    "→", font_size=30, color=ACCENT_BLUE, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Nota
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(1.0),
                "En ISELIN: los datos salen de archivos Excel (operativos, choferes, rutas, siniestros),\n"
                "se limpian en Power Query y quedan disponibles en Power BI para crear visualizaciones.",
                font_size=14, color=SUBTLE_GRAY)


def slide_tablas(prs):
    """Slide 6: Tablas de datos principales vs. referencia"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Tipos de tablas en Power BI",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    # Tabla principal (izquierda)
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
             "Tablas de datos principales",
             "Contienen los registros de eventos y transacciones.\n"
             "Cada fila es algo que pasó.\n\n"
             "Ejemplos en ISELIN:\n"
             "  •  Cada viaje realizado\n"
             "  •  Cada boleto vendido\n"
             "  •  Cada siniestro registrado\n"
             "  •  Cada gasto operativo",
             accent=ACCENT_BLUE, title_size=18, body_size=15)

    # Tabla referencia (derecha)
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5),
             "Tablas de referencia",
             "Contienen datos descriptivos que no cambian seguido.\n"
             "Sirven para clasificar y agrupar.\n\n"
             "Ejemplos en ISELIN:\n"
             "  •  Lista de choferes\n"
             "  •  Lista de rutas\n"
             "  •  Lista de colectivos\n"
             "  •  Calendario de fechas",
             accent=ACCENT_TEAL, title_size=18, body_size=15)

    # Conexión visual
    add_textbox(slide, Inches(5.6), Inches(3.5), Inches(2.2), Inches(0.5),
                "se conectan\npor campos en\ncomún",
                font_size=12, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)


def slide_filtros(prs):
    """Slide 7: Cómo funcionan los filtros"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "¿Cómo funcionan los filtros en PBI?",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.8),
                "Cuando aplicamos un filtro, Power BI recalcula todos los valores automáticamente.\n"
                "No hace falta cambiar fórmulas ni tocar los datos.",
                font_size=18, color=LIGHT_GRAY)

    # Ejemplo visual: Sin filtro vs. Con filtro
    add_card(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.5),
             "Sin filtro",
             "Vemos los indicadores de TODA la empresa:\n\n"
             "  •  Total viajes: 12.500\n"
             "  •  Km recorridos: 890.000\n"
             "  •  Siniestros: 45\n"
             "  •  Consumo gasoil: 320.000 lts",
             accent=MID_GRAY, title_size=18, body_size=15)

    add_card(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(3.5),
             "Filtro: Ruta 5",
             "Vemos solo los indicadores de la Ruta 5:\n\n"
             "  •  Total viajes: 1.800\n"
             "  •  Km recorridos: 125.000\n"
             "  •  Siniestros: 6\n"
             "  •  Consumo gasoil: 47.000 lts",
             accent=ACCENT_TEAL, title_size=18, body_size=15)

    # Flecha
    add_textbox(slide, Inches(5.8), Inches(3.8), Inches(1.5), Inches(0.5),
                "filtrar →", font_size=16, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)


def slide_estadisticas(prs):
    """Slide 8: Estadísticas básicas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Estadísticas básicas",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                "Tres números clave para entender los datos de ISELIN",
                font_size=18, color=LIGHT_GRAY)

    stats = [
        ("Promedio", "El valor \"medio\" de un conjunto de datos",
         "Km promedio recorridos\npor colectivo", ACCENT_BLUE),
        ("Mediana", "El valor del medio cuando ordenamos los datos.\nNo se distorsiona por valores extremos",
         "Tiempo mediano de demora\nen una ruta", ACCENT_TEAL),
        ("Dispersión", "Cuánto varían los datos respecto al promedio.\nMucha dispersión = datos muy diferentes entre sí",
         "Variabilidad del consumo\nde gasoil entre colectivos", ORANGE_ACCENT),
    ]

    for i, (titulo, definicion, ejemplo, color) in enumerate(stats):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(2.0)

        card = add_shape(slide, x, y, Inches(3.8), Inches(4.2), CARD_BG,
                         MSO_SHAPE.ROUNDED_RECTANGLE)

        # Acento
        add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Pt(3), color)

        # Título
        add_textbox(slide, x + Inches(0.2), y + Inches(0.3),
                    Inches(3.4), Inches(0.5),
                    titulo, font_size=22, color=color, bold=True)

        # Definición
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85),
                    Inches(3.4), Inches(1.4),
                    definicion, font_size=14, color=WHITE)

        # Ejemplo ISELIN
        add_shape(slide, x + Inches(0.15), y + Inches(2.5),
                  Inches(3.5), Pt(1), RGBColor(0x2A, 0x3A, 0x4A))
        add_textbox(slide, x + Inches(0.2), y + Inches(2.65),
                    Inches(3.4), Inches(0.3),
                    "En ISELIN:", font_size=12, color=SUBTLE_GRAY, bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(2.95),
                    Inches(3.4), Inches(1.0),
                    ejemplo, font_size=14, color=LIGHT_GRAY)


def slide_demo_pbi(prs):
    """Slide 9: Power BI - Demo en vivo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Power BI: conociendo la herramienta",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.25), Inches(5), Inches(0.5),
                "Demo en vivo  ·  35 minutos",
                font_size=16, color=ACCENT_TEAL, bold=True)

    # Las 3 vistas
    add_textbox(slide, Inches(0.8), Inches(1.85), Inches(5), Inches(0.5),
                "Las 3 vistas de PBI Desktop:", font_size=18, color=WHITE, bold=True)

    vistas = [
        ("Informe", "Donde se crean las\nvisualizaciones y dashboards"),
        ("Datos", "Donde se ven las tablas\ny los datos cargados"),
        ("Modelo", "Donde se conectan\nlas tablas entre sí"),
    ]

    for i, (titulo, desc) in enumerate(vistas):
        x = Inches(0.8) + Inches(i * 2.6)
        y = Inches(2.5)
        color = [ACCENT_BLUE, ACCENT_TEAL, GREEN_ACCENT][i]
        add_number_badge(slide, x, y, str(i + 1), color)
        add_textbox(slide, x + Inches(0.65), y + Inches(0.05),
                    Inches(1.8), Inches(0.35),
                    titulo, font_size=16, color=color, bold=True)
        add_textbox(slide, x, y + Inches(0.6), Inches(2.4), Inches(0.9),
                    desc, font_size=13, color=LIGHT_GRAY)

    # Práctica en vivo (lado derecho)
    add_card(slide, Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.8),
             "Práctica en vivo",
             "1. Cargar un archivo Excel real de ISELIN\n\n"
             "2. Vista previa en Power Query\n\n"
             "3. Ajustar tipos de datos\n\n"
             "4. Cerrar y aplicar\n\n"
             "5. Crear el primer gráfico de barras\n   con un indicador real",
             accent=ACCENT_TEAL, title_size=16, body_size=14)

    # Placeholder para captura
    placeholder = add_shape(slide, Inches(0.8), Inches(4.2),
                            Inches(6.8), Inches(2.8),
                            RGBColor(0x14, 0x20, 0x35),
                            MSO_SHAPE.ROUNDED_RECTANGLE)

    add_textbox(slide, Inches(2.5), Inches(5.2), Inches(3.5), Inches(0.5),
                "[ Espacio para captura de pantalla ]",
                font_size=14, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)


def slide_interfaz(prs):
    """Slide 10: Interfaz de PBI Desktop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Interfaz de Power BI Desktop",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    # Elementos de la interfaz
    elementos = [
        ("Cinta de opciones", "Menú principal con todas las\nherramientas organizadas\npor pestañas"),
        ("Panel de campos", "Lista de tablas y columnas\ndisponibles para usar en\nlas visualizaciones"),
        ("Panel de visualizaciones", "Tipos de gráficos disponibles\ny opciones de formato\npara cada uno"),
        ("Panel de filtros", "Filtros que aplican a todo\nel reporte, a una página\no a un gráfico específico"),
    ]

    for i, (titulo, desc) in enumerate(elementos):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + Inches(col * 5.8)
        y = Inches(1.5) + Inches(row * 2.3)

        add_card(slide, x, y, Inches(5.2), Inches(1.9),
                 titulo, desc,
                 accent=[ACCENT_BLUE, ACCENT_TEAL, ACCENT_TEAL, ACCENT_BLUE][i],
                 title_size=17, body_size=14)

    # Placeholder para captura
    placeholder = add_shape(slide, Inches(0.8), Inches(5.7),
                            Inches(11.8), Inches(1.3),
                            RGBColor(0x14, 0x20, 0x35),
                            MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(4.0), Inches(6.0), Inches(5.5), Inches(0.5),
                "[ Espacio para captura de la interfaz de PBI Desktop ]",
                font_size=14, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)


def slide_roadmap(prs):
    """Slide 11: Mapa de ruta - los 4 encuentros"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Mapa de ruta: los 4 encuentros",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    encuentros = [
        ("E1 — Hoy", "16/04", "Fundamentos,\nCalidad y\nPrimer Contacto", ACCENT_BLUE),
        ("E2", "21/04", "Modelo de Datos,\nPower Query y\nVisualizaciones", ACCENT_TEAL),
        ("E3", "23/04", "Fórmulas DAX,\nDashboards y\nStorytelling", ACCENT_BLUE),
        ("E4", "30/04", "Interpretación,\nPublicación y\nActualización", ACCENT_TEAL),
    ]

    # Línea de tiempo
    add_shape(slide, Inches(1.5), Inches(3.25), Inches(10.5), Pt(2), RGBColor(0x2A, 0x3A, 0x4A))

    for i, (titulo, fecha, desc, color) in enumerate(encuentros):
        x = Inches(1.2) + Inches(i * 3.05)
        y = Inches(1.6)

        # Punto en la línea
        add_shape(slide, x + Inches(1.05), Inches(3.12), Inches(0.3), Inches(0.3),
                  color, MSO_SHAPE.OVAL)

        # Título arriba
        add_textbox(slide, x, y, Inches(2.7), Inches(0.5),
                    titulo, font_size=18, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Fecha
        add_textbox(slide, x, y + Inches(0.45), Inches(2.7), Inches(0.35),
                    fecha, font_size=14, color=SUBTLE_GRAY,
                    alignment=PP_ALIGN.CENTER)

        # Descripción debajo de la línea
        add_textbox(slide, x, Inches(3.7), Inches(2.7), Inches(1.5),
                    desc, font_size=14, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Highlight E1
    highlight = add_shape(slide, Inches(0.9), Inches(1.4),
                          Inches(3.2), Inches(4.1),
                          RGBColor(0x34, 0x98, 0xDB))
    highlight.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x40)

    # Mover highlight al fondo
    slide.shapes._spTree.remove(highlight._element)
    slide.shapes._spTree.insert(2, highlight._element)


def slide_tarea(prs):
    """Slide 12: Tarea para el próximo encuentro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Tarea para el Encuentro 2",
                font_size=34, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(0.8))

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.6),
                "Para practicar antes del martes 21/04",
                font_size=18, color=LIGHT_GRAY)

    # Tarea principal
    add_card(slide, Inches(0.8), Inches(2.2), Inches(7.5), Inches(3.0),
             "Abrir Power BI Desktop y cargar un Excel de tu área",
             "1.  Abrir Power BI Desktop en tu computadora\n\n"
             "2.  Ir a Inicio → Obtener datos → Excel\n\n"
             "3.  Seleccionar un archivo Excel de tu área\n\n"
             "4.  Revisar la vista previa y hacer clic en \"Cargar\"\n\n"
             "5.  Explorar los datos en la vista de Datos",
             accent=ACCENT_TEAL, title_size=20, body_size=16)

    # Nota
    add_card(slide, Inches(9.0), Inches(2.2), Inches(3.8), Inches(3.0),
             "Si tenés dudas",
             "No te preocupes si algo\nno funciona. Arrancamos\nel Encuentro 2 resolviendo\ndudas y errores comunes.\n\nLo importante es intentarlo.",
             accent=MID_GRAY, title_size=16, body_size=14)

    # Entregables
    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(5), Inches(0.5),
                "Entregables que van a recibir:", font_size=18, color=WHITE, bold=True)

    entregables = [
        "Archivo .pbix base con la primera tabla cargada",
        "Guía de conceptos fundamentales (PDF)",
        "Guía de instalación de PBI Desktop (PDF)",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(10), Inches(1.2),
                    entregables, font_size=15, color=LIGHT_GRAY)


def slide_cierre(prs):
    """Slide 13: Cierre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Barra lateral
    add_shape(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, ACCENT_BLUE)

    # Logo
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "data "
    run1.font.size = Pt(22)
    run1.font.color.rgb = MID_GRAY
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "oilers"
    run2.font.size = Pt(22)
    run2.font.color.rgb = ACCENT_TEAL
    run2.font.bold = True
    run2.font.name = "Calibri"

    # Línea de acento
    add_accent_line(slide, Inches(0.6), Inches(3.0), Inches(1.2))

    # Texto de cierre
    add_textbox(slide, Inches(0.6), Inches(3.3), Inches(10), Inches(1),
                "¡Gracias!",
                font_size=48, color=WHITE, bold=True)

    add_textbox(slide, Inches(0.6), Inches(4.3), Inches(10), Inches(0.6),
                "Nos vemos en el Encuentro 2 — Martes 21 de abril",
                font_size=22, color=LIGHT_GRAY)

    add_textbox(slide, Inches(0.6), Inches(5.0), Inches(10), Inches(0.6),
                "Modelo de Datos, Power Query y Primeras Visualizaciones",
                font_size=18, color=ACCENT_TEAL)

    # Contacto / preguntas
    add_textbox(slide, Inches(0.6), Inches(6.2), Inches(10), Inches(0.5),
                "¿Preguntas? Escribinos por el canal de Teams/WhatsApp del programa",
                font_size=14, color=SUBTLE_GRAY)


# ─── Main ───

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_portada(prs)
    slide_agenda(prs)
    slide_por_que_importa(prs)
    slide_ciclo_mejora(prs)
    slide_etl(prs)
    slide_tablas(prs)
    slide_filtros(prs)
    slide_estadisticas(prs)
    slide_demo_pbi(prs)
    slide_interfaz(prs)
    slide_roadmap(prs)
    slide_tarea(prs)
    slide_cierre(prs)

    output = "presentacion_E1.pptx"
    prs.save(output)
    print(f"Presentación generada: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
