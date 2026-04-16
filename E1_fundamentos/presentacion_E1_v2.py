#!/usr/bin/env python3
"""
Presentación E1 v2 - Estilo claro, moderno, profesional.
Tema: light suave (no agresivo) con acentos DATAOILERS.

Uso: python3 presentacion_E1_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Dimensiones ───
SW = Inches(13.333)
SH = Inches(7.5)

# ─── Paleta ───
BG = RGBColor(0xF0, 0xF2, 0xF5)          # Fondo general
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # Tarjetas
BG_SOFT = RGBColor(0xE8, 0xEB, 0xF0)     # Cards secundarias
BORDER = RGBColor(0xDE, 0xE2, 0xE8)      # Bordes suaves
DARK = RGBColor(0x1A, 0x2A, 0x3A)        # Títulos
BODY = RGBColor(0x4A, 0x55, 0x68)        # Texto cuerpo
SUBTLE = RGBColor(0x8A, 0x96, 0xA6)      # Texto secundario
BLUE = RGBColor(0x34, 0x98, 0xDB)        # Acento principal
TEAL = RGBColor(0x4A, 0xAD, 0xA4)        # Acento secundario
NAVY = RGBColor(0x0F, 0x1B, 0x2D)        # Fondo oscuro (portada/cierre)
NAVY2 = RGBColor(0x16, 0x23, 0x3E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED_SOFT = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0x8F, 0xA8, 0xBF)
SGRAY = RGBColor(0x5A, 0x7A, 0x94)

# Fondos de tarjeta para tema claro
CARD_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
CARD_TEAL = RGBColor(0xE8, 0xF8, 0xF5)
CARD_ORANGE = RGBColor(0xFD, 0xF2, 0xE9)
CARD_GREEN = RGBColor(0xEA, 0xFB, 0xF0)


# ─── Helpers ───

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, color, rounded=False):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def rect_border(slide, l, t, w, h, fill_color, border_color, rounded=True):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color
    s.line.width = Pt(1)
    return s


def txt(slide, l, t, w, h, text, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb


def multitext(slide, l, t, w, h, lines, default_size=16, default_color=BODY):
    """lines = [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        size = item[1] if len(item) > 1 else default_size
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def accent_line(slide, l, t, w, color=BLUE):
    return rect(slide, l, t, w, Pt(4), color)


def badge(slide, l, t, num, color=BLUE, text_color=WHITE):
    rect(slide, l, t, Inches(0.5), Inches(0.5), color, rounded=True)
    txt(slide, l, t + Inches(0.03), Inches(0.5), Inches(0.44),
        str(num), size=18, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, title, body, accent=BLUE, bg_color=BG_WHITE,
         title_size=15, body_size=13, border=True):
    if border:
        rect_border(slide, l, t, w, h, bg_color, BORDER)
    else:
        rect(slide, l, t, w, h, bg_color, rounded=True)
    # Accent
    rect(slide, l + Inches(0.15), t + Inches(0.12), Inches(0.45), Pt(3), accent)
    # Title
    txt(slide, l + Inches(0.18), t + Inches(0.25), w - Inches(0.35), Inches(0.35),
        title, size=title_size, color=DARK, bold=True)
    # Body
    txt(slide, l + Inches(0.18), t + Inches(0.6), w - Inches(0.35), h - Inches(0.7),
        body, size=body_size, color=BODY)


def wiki_callout(slide, l, t, text="Profundizamos en la wiki durante el encuentro"):
    rect_border(slide, l, t, Inches(4.5), Inches(0.4), CARD_BLUE, BLUE, rounded=True)
    txt(slide, l + Inches(0.15), t + Inches(0.05), Inches(4.2), Inches(0.3),
        text, size=11, color=BLUE, bold=False)


def section_header(slide, number, title, subtitle=""):
    """Header estándar para slides de contenido."""
    # Barra superior azul muy sutil
    rect(slide, Inches(0), Inches(0), SW, Inches(0.06), BLUE)

    # Número + título
    txt(slide, Inches(0.8), Inches(0.35), Inches(0.5), Inches(0.5),
        f"{number:02d}", size=14, color=BLUE, bold=True)
    txt(slide, Inches(1.25), Inches(0.3), Inches(10), Inches(0.6),
        title, size=30, color=DARK, bold=True)
    accent_line(slide, Inches(1.25), Inches(0.9), Inches(0.7))

    if subtitle:
        txt(slide, Inches(1.25), Inches(1.05), Inches(10), Inches(0.5),
            subtitle, size=16, color=SUBTLE)


# ─── SLIDES ───

def s_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    # Sidebar
    rect(slide, Inches(0), Inches(0), Inches(0.1), SH, BLUE)

    # Logo
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(4), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "data "
    r1.font.size = Pt(20)
    r1.font.color.rgb = LGRAY
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "oilers"
    r2.font.size = Pt(20)
    r2.font.color.rgb = TEAL
    r2.font.bold = True
    r2.font.name = "Calibri"

    # Encounter badge
    rect(slide, Inches(0.55), Inches(2.8), Inches(1.4), Inches(0.4), BLUE, rounded=True)
    txt(slide, Inches(0.55), Inches(2.82), Inches(1.4), Inches(0.36),
        "ENCUENTRO 1", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    accent_line(slide, Inches(0.55), Inches(3.5), Inches(1))

    txt(slide, Inches(0.55), Inches(3.7), Inches(8), Inches(1.5),
        "Fundamentos, Calidad\ny Primer Contacto con\nPower BI",
        size=40, color=WHITE, bold=True)

    txt(slide, Inches(0.55), Inches(5.5), Inches(8), Inches(0.5),
        "Capacitación en Análisis de Datos · ISELIN",
        size=16, color=LGRAY)

    txt(slide, Inches(0.55), Inches(6.5), Inches(6), Inches(0.4),
        "Jueves 16 de abril de 2026  ·  Presencial  ·  2 horas",
        size=13, color=SGRAY)


def s_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 1, "Agenda del día",
                   "Lo que vamos a recorrer en estas 2 horas")

    bloques = [
        ("10 min", "Apertura y contexto", "Presentación del programa y expectativas", BLUE),
        ("25 min", "Análisis de datos y calidad", "Por qué importa y cómo se conecta con ISELIN", TEAL),
        ("25 min", "Conceptos fundamentales", "ETL, tablas, filtros y estadísticas básicas", BLUE),
        ("35 min", "Power BI: demo en vivo", "Conocer la herramienta y cargar datos reales", TEAL),
        ("15 min", "Cierre y tarea", "Mapa de ruta y preparación para el Encuentro 2", BLUE),
        ("10 min", "Tiempo extra", "Preguntas adicionales", SUBTLE),
    ]

    for i, (tiempo, tema, desc, color) in enumerate(bloques):
        y = Inches(1.55) + Inches(i * 0.9)
        # Card
        rect_border(slide, Inches(1.0), y, Inches(11.5), Inches(0.75), BG_WHITE, BORDER)
        # Accent bar
        rect(slide, Inches(1.0), y, Pt(5), Inches(0.75), color)
        # Number
        txt(slide, Inches(1.2), y + Inches(0.15), Inches(0.4), Inches(0.45),
            str(i + 1) if i < 5 else "—", size=18, color=color, bold=True)
        # Tiempo
        txt(slide, Inches(1.7), y + Inches(0.18), Inches(1.2), Inches(0.4),
            tiempo, size=13, color=SUBTLE, bold=True)
        # Tema
        txt(slide, Inches(3.0), y + Inches(0.1), Inches(4), Inches(0.3),
            tema, size=17, color=DARK, bold=True)
        # Desc
        txt(slide, Inches(3.0), y + Inches(0.4), Inches(5), Inches(0.3),
            desc, size=12, color=SUBTLE)
        # Barra de proporción
        max_min = 35
        mins = int(tiempo.split()[0])
        bar_w = Inches(2.5 * mins / max_min)
        rect(slide, Inches(9.5), y + Inches(0.28), bar_w, Inches(0.2), color, rounded=True)


def s_por_que(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 2, "¿Por qué importa el análisis de datos?",
                   "De la intuición a las decisiones basadas en evidencia")

    cards_data = [
        ("Decisiones a ciegas",
         "El 33% de los ejecutivos aún basan sus decisiones\nen intuición personal (PwC, 2016).\n\n"
         "La intuición es valiosa, pero insuficiente\nen un entorno de alta complejidad.",
         CARD_ORANGE, ORANGE),
        ("Datos como activo",
         "Las empresas que aprovechan sus datos\naventajan a sus competidores en un 85%\n"
         "en crecimiento de ventas.\n\n"
         "Los datos no son un costo de IT,\nson un activo generador de valor.",
         CARD_GREEN, GREEN),
        ("Visibilidad compartida",
         "Un dashboard permite que todas las áreas\nvean los mismos números.\n\n"
         "Alinea criterios, facilita la comunicación\ny reduce las discusiones\nbasadas en suposiciones.",
         CARD_BLUE, BLUE),
    ]

    for i, (title, body, bg_c, accent) in enumerate(cards_data):
        x = Inches(0.8) + Inches(i * 4.15)
        card(slide, x, Inches(1.8), Inches(3.85), Inches(3.8),
             title, body, accent=accent, bg_color=bg_c,
             title_size=16, body_size=13)

    wiki_callout(slide, Inches(0.8), Inches(6.0),
                 "Wiki Cap. 1: La Evolución de la Intuición al Dato")


def s_piramide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 3, "De datos a decisiones",
                   "La pirámide del conocimiento (DIKW)")

    niveles = [
        ("DATO", "Materia prima. Registros brutos\nsin contexto ni significado.",
         "Un registro de viaje:\nfecha, chofer, km, gasoil", BG_SOFT, SUBTLE,
         Inches(8.5)),
        ("INFORMACIÓN", "Datos organizados que responden\npreguntas básicas: qué, cuándo, cuánto.",
         "Total de viajes del mes: 1.250\nKm promedio por día: 4.500", CARD_BLUE, BLUE,
         Inches(6.8)),
        ("CONOCIMIENTO", "Información interpretada.\nEntendemos patrones y causas.",
         "Los viajes bajaron 15% porque\nhubo 3 colectivos en reparación", CARD_TEAL, TEAL,
         Inches(5.1)),
        ("SABIDURÍA", "Conocimiento aplicado para\ntomar decisiones y actuar.",
         "Programar mantenimiento preventivo\npara evitar la baja de flota", CARD_GREEN, GREEN,
         Inches(3.4)),
    ]

    for i, (titulo, desc, ejemplo, bg_c, color, w) in enumerate(niveles):
        y = Inches(1.5) + Inches(i * 1.35)
        x = Inches((13.333 - w.inches) / 2)  # Centrado

        rect_border(slide, x, y, w, Inches(1.15), bg_c, BORDER)
        # Título
        txt(slide, x + Inches(0.2), y + Inches(0.08), Inches(2), Inches(0.3),
            titulo, size=13, color=color, bold=True)
        # Desc
        txt(slide, x + Inches(0.2), y + Inches(0.35), Inches(3.5), Inches(0.7),
            desc, size=11, color=BODY)
        # Ejemplo
        txt(slide, x + w - Inches(4), y + Inches(0.15), Inches(3.7), Inches(0.85),
            ejemplo, size=11, color=SUBTLE)
        # Flecha
        if i < 3:
            txt(slide, Inches(6.3), y + Inches(1.05), Inches(0.8), Inches(0.3),
                "▲", size=14, color=color, align=PP_ALIGN.CENTER)

    wiki_callout(slide, Inches(0.8), Inches(6.95),
                 "Wiki Cap. 2: Alfabetización de Datos — Pirámide DIKW")


def s_niveles_analitica(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 4, "Los 4 niveles de análisis",
                   "Cada nivel responde una pregunta diferente sobre el negocio")

    niveles = [
        ("Descriptivo", "¿Qué pasó?",
         "Reportes, dashboards,\nresúmenes de datos\nhistóricos",
         "Donde arrancamos\nen esta capacitación", BLUE, CARD_BLUE),
        ("Diagnóstico", "¿Por qué pasó?",
         "Análisis de causas,\ncorrelaciones,\ndrill-down en reportes",
         "Encuentros 2 y 3", TEAL, CARD_TEAL),
        ("Predictivo", "¿Qué va a pasar?",
         "Modelos estadísticos,\ntendencias,\nforecasting",
         "Encuentro 4\n(introducción)", ORANGE, CARD_ORANGE),
        ("Prescriptivo", "¿Qué debemos hacer?",
         "Recomendaciones\nautomáticas,\noptimización",
         "Futuro", GREEN, CARD_GREEN),
    ]

    for i, (titulo, pregunta, desc, cuando, color, bg_c) in enumerate(niveles):
        x = Inches(0.6) + Inches(i * 3.15)
        y = Inches(1.7)

        rect_border(slide, x, y, Inches(2.95), Inches(4.5), bg_c, BORDER)

        # Número
        badge(slide, x + Inches(0.15), y + Inches(0.15), i + 1, color)

        # Título
        txt(slide, x + Inches(0.75), y + Inches(0.2), Inches(2), Inches(0.35),
            titulo, size=17, color=color, bold=True)

        # Pregunta
        txt(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.6), Inches(0.4),
            pregunta, size=20, color=DARK, bold=True)

        # Descripción
        txt(slide, x + Inches(0.15), y + Inches(1.3), Inches(2.6), Inches(1.5),
            desc, size=13, color=BODY)

        # Cuándo
        rect(slide, x + Inches(0.15), y + Inches(3.2), Inches(2.6), Pt(1), BORDER)
        txt(slide, x + Inches(0.15), y + Inches(3.35), Inches(2.6), Inches(0.6),
            cuando, size=11, color=SUBTLE)

    # Flecha de complejidad
    rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Pt(2), BORDER)
    txt(slide, Inches(0.6), Inches(6.6), Inches(3), Inches(0.3),
        "Menor complejidad", size=11, color=SUBTLE)
    txt(slide, Inches(9.7), Inches(6.6), Inches(3), Inches(0.3),
        "Mayor valor para el negocio →", size=11, color=SUBTLE, align=PP_ALIGN.RIGHT)

    wiki_callout(slide, Inches(4.5), Inches(6.55),
                 "Wiki Cap. 5: Analítica para la Acción")


def s_calidad(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 5, "Calidad de datos: la base de todo",
                   "\"Basura entra, basura sale\" — Si los datos son malos, las decisiones también")

    dims = [
        ("Precisión", "Los datos reflejan\nla realidad", BLUE),
        ("Completitud", "No faltan valores\nimportantes", TEAL),
        ("Consistencia", "Los mismos datos\nno se contradicen", GREEN),
        ("Validez", "Formatos y rangos\ncorrectos", ORANGE),
        ("Unicidad", "Sin registros\nduplicados", BLUE),
        ("Oportunidad", "Datos actualizados\ncuando se necesitan", TEAL),
    ]

    for i, (titulo, desc, color) in enumerate(dims):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + Inches(col * 4.15)
        y = Inches(1.7) + Inches(row * 2.2)

        rect_border(slide, x, y, Inches(3.85), Inches(1.8), BG_WHITE, BORDER)
        rect(slide, x, y, Pt(5), Inches(1.8), color)
        badge(slide, x + Inches(0.2), y + Inches(0.2), i + 1, color)
        txt(slide, x + Inches(0.85), y + Inches(0.22), Inches(2.8), Inches(0.35),
            titulo, size=17, color=DARK, bold=True)
        txt(slide, x + Inches(0.85), y + Inches(0.6), Inches(2.8), Inches(1.0),
            desc, size=14, color=BODY)

    wiki_callout(slide, Inches(0.8), Inches(6.3),
                 "Wiki Cap. 4: Calidad y Gobierno de Datos")


def s_ciclo_mejora(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 6, "Ciclo de mejora continua",
                   "Power BI es la herramienta de \"Verificar\" dentro del ciclo")

    etapas = [
        ("Planificar", "Definir objetivos\ny planificar acciones", BLUE),
        ("Hacer", "Ejecutar lo\nplanificado", TEAL),
        ("Verificar", "Medir resultados\ncon datos reales", ORANGE),
        ("Actuar", "Ajustar y mejorar\nel proceso", GREEN),
    ]

    for i, (titulo, desc, color) in enumerate(etapas):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(2.0)

        is_verificar = (i == 2)
        bg_c = CARD_ORANGE if is_verificar else BG_WHITE

        rect_border(slide, x, y, Inches(2.85), Inches(2.8), bg_c, color if is_verificar else BORDER)
        badge(slide, x + Inches(1.1), y + Inches(0.2), i + 1, color)
        txt(slide, x, y + Inches(0.85), Inches(2.85), Inches(0.4),
            titulo, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, x, y + Inches(1.35), Inches(2.85), Inches(1.0),
            desc, size=14, color=BODY, align=PP_ALIGN.CENTER)

        if is_verificar:
            txt(slide, x, y + Inches(2.2), Inches(2.85), Inches(0.4),
                "Power BI", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

        # Flechas
        if i < 3:
            txt(slide, x + Inches(2.65), y + Inches(1.1), Inches(0.7), Inches(0.5),
                "→", size=26, color=SUBTLE, align=PP_ALIGN.CENTER)

    txt(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.8),
        "ISELIN ya trabaja con mejora continua.\n"
        "Power BI se suma como la herramienta que permite verificar con datos reales si las acciones dan resultado.",
        size=15, color=BODY)


def s_etl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 7, "Del Excel a Power BI",
                   "El proceso de carga en 3 pasos")

    pasos = [
        ("Extraer", "Sacar los datos desde\nExcel, Wara u otra fuente",
         "Los archivos Excel de ISELIN\nson el punto de partida", BLUE, CARD_BLUE),
        ("Transformar", "Limpiar, corregir formatos,\neliminar datos innecesarios",
         "Power Query es el \"taller mecánico\"\ndonde se preparan los datos", TEAL, CARD_TEAL),
        ("Cargar", "Meter los datos limpios en\nPower BI, listos para usar",
         "Los datos quedan disponibles\npara crear visualizaciones", GREEN, CARD_GREEN),
    ]

    for i, (titulo, desc, ejemplo, color, bg_c) in enumerate(pasos):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(1.7)

        rect_border(slide, x, y, Inches(3.85), Inches(3.8), bg_c, BORDER)
        badge(slide, x + Inches(0.15), y + Inches(0.15), i + 1, color)
        txt(slide, x + Inches(0.8), y + Inches(0.2), Inches(2.8), Inches(0.35),
            titulo, size=22, color=color, bold=True)

        txt(slide, x + Inches(0.2), y + Inches(0.75), Inches(3.4), Inches(1.2),
            desc, size=15, color=DARK)

        rect(slide, x + Inches(0.2), y + Inches(2.1), Inches(3.4), Pt(1), BORDER)
        txt(slide, x + Inches(0.2), y + Inches(2.3), Inches(3.4), Inches(0.3),
            "En ISELIN:", size=12, color=SUBTLE, bold=True)
        txt(slide, x + Inches(0.2), y + Inches(2.6), Inches(3.4), Inches(1.0),
            ejemplo, size=13, color=BODY)

        if i < 2:
            txt(slide, x + Inches(3.65), y + Inches(1.5), Inches(0.7), Inches(0.5),
                "→", size=28, color=color, bold=True, align=PP_ALIGN.CENTER)

    wiki_callout(slide, Inches(0.8), Inches(5.9),
                 "Wiki Cap. 3: La Fábrica de Datos — ETL y ELT")


def s_tablas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 8, "Tipos de tablas en Power BI",
                   "Los datos se organizan en dos tipos de tablas que se conectan entre sí")

    card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.2),
         "Tablas de datos principales",
         "Contienen los registros de eventos.\nCada fila es algo que pasó.\n\n"
         "Ejemplos en ISELIN:\n"
         "  ·  Cada viaje realizado\n"
         "  ·  Cada boleto vendido\n"
         "  ·  Cada siniestro registrado\n"
         "  ·  Cada gasto operativo",
         accent=BLUE, bg_color=CARD_BLUE, title_size=18, body_size=14)

    card(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.2),
         "Tablas de referencia",
         "Contienen datos descriptivos que\nno cambian seguido. Sirven para\nclasificar y agrupar.\n\n"
         "Ejemplos en ISELIN:\n"
         "  ·  Lista de choferes\n"
         "  ·  Lista de rutas\n"
         "  ·  Lista de colectivos\n"
         "  ·  Calendario de fechas",
         accent=TEAL, bg_color=CARD_TEAL, title_size=18, body_size=14)

    # Conexión
    rect_border(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.9), BG_WHITE, BORDER)
    txt(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.7),
        "Se conectan por\ncampos en común\n(ej: ID de chofer)",
        size=12, color=SUBTLE, align=PP_ALIGN.CENTER)


def s_filtros(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 9, "¿Cómo funcionan los filtros?",
                   "Al aplicar un filtro, Power BI recalcula todo automáticamente")

    card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.8),
         "Sin filtro — Toda la empresa",
         "Total viajes:            12.500\n"
         "Km recorridos:       890.000\n"
         "Siniestros:                    45\n"
         "Consumo gasoil:    320.000 lts",
         accent=SUBTLE, bg_color=BG_WHITE, title_size=18, body_size=16)

    card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.8),
         "Filtro aplicado: Ruta 5",
         "Total viajes:              1.800\n"
         "Km recorridos:       125.000\n"
         "Siniestros:                      6\n"
         "Consumo gasoil:      47.000 lts",
         accent=TEAL, bg_color=CARD_TEAL, title_size=18, body_size=16)

    txt(slide, Inches(5.6), Inches(3.3), Inches(2), Inches(0.5),
        "filtrar →", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    txt(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.6),
        "No hace falta cambiar fórmulas ni tocar los datos. PBI hace todo el trabajo.",
        size=15, color=BODY)


def s_estadisticas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 10, "Estadísticas básicas",
                   "Tres números clave para entender los datos de ISELIN")

    stats = [
        ("Promedio", "El valor \"medio\" de un\nconjunto de datos.",
         "Km promedio recorridos\npor colectivo", BLUE, CARD_BLUE),
        ("Mediana", "El valor del medio cuando\nordenamos los datos.\nNo se distorsiona por extremos.",
         "Tiempo mediano de\ndemora en una ruta", TEAL, CARD_TEAL),
        ("Dispersión", "Cuánto varían los datos.\nMucha dispersión = datos\nmuy diferentes entre sí.",
         "Variabilidad del consumo\nde gasoil entre colectivos", ORANGE, CARD_ORANGE),
    ]

    for i, (titulo, desc, ejemplo, color, bg_c) in enumerate(stats):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(1.7)

        rect_border(slide, x, y, Inches(3.85), Inches(4.2), bg_c, BORDER)
        rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Pt(3), color)
        txt(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.45),
            titulo, size=22, color=color, bold=True)
        txt(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.4), Inches(1.5),
            desc, size=14, color=DARK)

        rect(slide, x + Inches(0.2), y + Inches(2.5), Inches(3.4), Pt(1), BORDER)
        txt(slide, x + Inches(0.2), y + Inches(2.65), Inches(3.4), Inches(0.3),
            "En ISELIN:", size=12, color=SUBTLE, bold=True)
        txt(slide, x + Inches(0.2), y + Inches(2.95), Inches(3.4), Inches(1.0),
            ejemplo, size=14, color=BODY)


def s_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 11, "Power BI: demo en vivo",
                   "35 minutos — Conocer la herramienta y cargar datos reales")

    # Las 3 vistas
    vistas = [
        ("Informe", "Donde se crean\nlas visualizaciones", BLUE),
        ("Datos", "Donde se ven\nlas tablas cargadas", TEAL),
        ("Modelo", "Donde se conectan\nlas tablas entre sí", GREEN),
    ]

    for i, (titulo, desc, color) in enumerate(vistas):
        x = Inches(0.8) + Inches(i * 2.6)
        y = Inches(1.6)
        rect_border(slide, x, y, Inches(2.35), Inches(1.6), BG_WHITE, BORDER)
        badge(slide, x + Inches(0.12), y + Inches(0.12), i + 1, color)
        txt(slide, x + Inches(0.7), y + Inches(0.15), Inches(1.5), Inches(0.3),
            titulo, size=16, color=color, bold=True)
        txt(slide, x + Inches(0.12), y + Inches(0.6), Inches(2.1), Inches(0.85),
            desc, size=13, color=BODY)

    # Pasos de la práctica
    card(slide, Inches(8.5), Inches(1.6), Inches(4.2), Inches(4.5),
         "Pasos de la práctica",
         "1.  Cargar un Excel real de ISELIN\n\n"
         "2.  Vista previa en Power Query\n\n"
         "3.  Ajustar tipos de datos\n\n"
         "4.  Cerrar y aplicar\n\n"
         "5.  Crear el primer gráfico de barras\n"
         "     con un indicador real de ISELIN",
         accent=TEAL, title_size=16, body_size=14)

    # Placeholder captura
    rect_border(slide, Inches(0.8), Inches(3.5), Inches(7.2), Inches(3.3), BG_SOFT, BORDER)
    txt(slide, Inches(2.5), Inches(4.8), Inches(4), Inches(0.5),
        "[ Espacio para captura de pantalla de PBI ]",
        size=14, color=SUBTLE, align=PP_ALIGN.CENTER)


def s_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 12, "Mapa de ruta: los 4 encuentros")

    encuentros = [
        ("E1 — Hoy", "16/04", "Fundamentos, Calidad\ny Primer Contacto", BLUE, True),
        ("E2", "21/04", "Modelo de Datos,\nPower Query y\nVisualizaciones", TEAL, False),
        ("E3", "23/04", "Fórmulas DAX,\nDashboards y\nStorytelling", BLUE, False),
        ("E4", "30/04", "Interpretación,\nPublicación y\nActualización", TEAL, False),
    ]

    # Timeline
    rect(slide, Inches(1.5), Inches(3.6), Inches(10.3), Pt(3), BORDER)

    for i, (titulo, fecha, desc, color, active) in enumerate(encuentros):
        x = Inches(1.0) + Inches(i * 3.05)

        # Dot
        dot_color = color if active else SUBTLE
        rect(slide, x + Inches(1.1), Inches(3.45), Inches(0.35), Inches(0.35),
             dot_color, rounded=True)

        # Card
        bg_c = CARD_BLUE if active else BG_WHITE
        rect_border(slide, x, Inches(1.5), Inches(2.7), Inches(1.6), bg_c,
                    color if active else BORDER)

        txt(slide, x, Inches(1.58), Inches(2.7), Inches(0.35),
            titulo, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, x, Inches(1.9), Inches(2.7), Inches(0.3),
            fecha, size=12, color=SUBTLE, align=PP_ALIGN.CENTER)

        # Desc below timeline
        txt(slide, x, Inches(4.2), Inches(2.7), Inches(1.5),
            desc, size=13, color=BODY, align=PP_ALIGN.CENTER)

    txt(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.6),
        "5 días entre E1 y E2  ·  2 días entre E2 y E3  ·  7 días entre E3 y E4  —  cada intervalo es oportunidad de práctica",
        size=13, color=SUBTLE)


def s_tarea(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BG)
    section_header(slide, 13, "Tarea para el Encuentro 2",
                   "Para practicar antes del martes 21/04")

    card(slide, Inches(0.8), Inches(1.7), Inches(7.5), Inches(3.2),
         "Abrir Power BI Desktop y cargar un Excel de tu área",
         "1.  Abrir Power BI Desktop\n\n"
         "2.  Ir a Inicio → Obtener datos → Excel\n\n"
         "3.  Seleccionar un archivo Excel de tu área\n\n"
         "4.  Revisar la vista previa y hacer clic en \"Cargar\"\n\n"
         "5.  Explorar los datos en la vista de Datos",
         accent=TEAL, title_size=18, body_size=15)

    card(slide, Inches(9.0), Inches(1.7), Inches(3.8), Inches(3.2),
         "Si tenés dudas",
         "No te preocupes si algo\nno funciona.\n\n"
         "Arrancamos el Encuentro 2\nresolviendo dudas y\nerrores comunes.\n\n"
         "Lo importante es intentarlo.",
         accent=SUBTLE, title_size=16, body_size=13)

    # Entregables
    txt(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.4),
        "Entregables que van a recibir:", size=18, color=DARK, bold=True)

    entregables = [
        ("Archivo .pbix base", "Con la primera tabla cargada y una visualización"),
        ("Guía de conceptos fundamentales", "ETL, tablas, filtros, estadísticas, calidad"),
        ("Guía de instalación de PBI Desktop", "Para que puedan practicar por su cuenta"),
    ]

    for i, (titulo, desc) in enumerate(entregables):
        y = Inches(5.8) + Inches(i * 0.5)
        badge(slide, Inches(0.8), y, i + 1, [BLUE, TEAL, GREEN][i])
        txt(slide, Inches(1.45), y + Inches(0.02), Inches(4), Inches(0.25),
            titulo, size=14, color=DARK, bold=True)
        txt(slide, Inches(5.5), y + Inches(0.04), Inches(5), Inches(0.25),
            desc, size=12, color=SUBTLE)


def s_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    rect(slide, Inches(0), Inches(0), Inches(0.1), SH, BLUE)

    # Logo
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(4), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "data "
    r1.font.size = Pt(20)
    r1.font.color.rgb = LGRAY
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "oilers"
    r2.font.size = Pt(20)
    r2.font.color.rgb = TEAL
    r2.font.bold = True
    r2.font.name = "Calibri"

    accent_line(slide, Inches(0.55), Inches(3.0), Inches(1))

    txt(slide, Inches(0.55), Inches(3.3), Inches(10), Inches(0.9),
        "¡Gracias!", size=48, color=WHITE, bold=True)
    txt(slide, Inches(0.55), Inches(4.3), Inches(10), Inches(0.5),
        "Nos vemos en el Encuentro 2 — Martes 21 de abril",
        size=22, color=LGRAY)
    txt(slide, Inches(0.55), Inches(4.9), Inches(10), Inches(0.5),
        "Modelo de Datos, Power Query y Primeras Visualizaciones",
        size=18, color=TEAL)
    txt(slide, Inches(0.55), Inches(6.3), Inches(10), Inches(0.4),
        "¿Preguntas? Escribinos por el canal de Teams/WhatsApp del programa",
        size=13, color=SGRAY)


# ─── MAIN ───

def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    s_portada(prs)       # 1
    s_agenda(prs)        # 2
    s_por_que(prs)       # 3
    s_piramide(prs)      # 4
    s_niveles_analitica(prs)  # 5
    s_calidad(prs)       # 6
    s_ciclo_mejora(prs)  # 7
    s_etl(prs)           # 8
    s_tablas(prs)        # 9
    s_filtros(prs)       # 10
    s_estadisticas(prs)  # 11
    s_demo(prs)          # 12
    s_roadmap(prs)       # 13
    s_tarea(prs)         # 14
    s_cierre(prs)        # 15

    out = "presentacion_E1_v2.pptx"
    prs.save(out)
    print(f"Presentación generada: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
